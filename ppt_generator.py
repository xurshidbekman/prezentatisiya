"""
====================================================
  SIMPLE PPTX GENERATOR — ppt_generator.py
  (Faqat matn va DuckDuckGo rasmlari)
====================================================
"""
import os
import json
import requests
import google.generativeai as genai
from dotenv import load_dotenv
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# DuckDuckGo dan rasm qidirish kutubxonasi
from duckduckgo_search import DDGS

load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

# ══════════════════════════════════════════════════
#  YORDAMCHI FUNKSIYALAR
# ══════════════════════════════════════════════════

def _set_white_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

def _text(slide, text, left, top, width, height,
          size=18, bold=False, italic=False,
          color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT):
    """Sodda matn qutisi yaratish usuli."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return tb

import time

import time
import urllib.parse

def _download_image(keyword: str, path: str) -> bool:
    """Rasmni yuklash: DDGS -> AI Generatsiya -> Wikipedia"""
    img_url = None
    
    # 1. DuckDuckGo (Haqiqiy rasmlar)
    try:
        time.sleep(1.2)
        results = DDGS().images(keyword, max_results=1)
        if results and len(results) > 0:
            img_url = results[0].get("image")
    except Exception as e:
        print(f"  [DDG xatosi/Limit] {e}. (Fallback ishga tushiriladi)")

    # 2. AI Rasm generatsiyasi (Pollinations) - Hech qanday limitsiz!
    if not img_url:
        try:
            # Pollinations AI yordamida tezkor chizish
            encoded_prompt = urllib.parse.quote(f"High quality educational illustration of {keyword}, clean background")
            img_url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=800&height=600&nologo=true"
        except Exception:
            pass

    if not img_url:
        return False
        
    # Yuklab olishga harakat (URL aniqlangandan keyin)
    try:
        data = requests.get(img_url, timeout=15).content
        with open(path, "wb") as f:
            f.write(data)
        return True
    except Exception as e:
        print(f"  [Yuklash xatosi] URL ishlamadi: {e}")
        return False

# ══════════════════════════════════════════════════
#  GEMINI — SLAYD KONTENT
# ══════════════════════════════════════════════════

def _get_slides_from_gemini(input_data: str, is_image: bool, slide_count: int) -> list:
    model = genai.GenerativeModel("gemini-3-flash-preview")

    prompt = f"""Siz aniq va londa ma'lumot beruvchi o'qituvchisiz.
Matn yoki rasm asosida aniq va tushunarli PPTX slaydlari tayyorlang (Aynan {slide_count} ta slayd bo'lsin).

FAQATGINA quyidagi JSON formatida javob bering:
[
  {{
    "title": "Sarlavha",
    "content": "Slayd haqida aniq, tushunarli ma'lumot matni",
    "bullets": ["Qisqacha element", "Ikkinchi fakt", "Uchinchi fakt"],
    "image_search_query": "An exact, visual English search query strictly representing the slide's content (e.g. 'Newton pendulum physics experiment', 'Solar system high quality')"
  }}
]
QOIDALAR:
- O'zbek tilida yozing. Hech qanday turli xil murakkab belgilar va formatlarsiz.
- image_search_query har doim albatta inglizcha va rasmlar aniq topilishi uchun juda spetsifik (aniq) bo'lsin.
"""

    if is_image:
        img = Image.open(input_data)
        response = model.generate_content([prompt, img])
    else:
        response = model.generate_content([prompt, input_data])

    raw = response.text.strip()
    for tag in ["```json", "```"]:
        if raw.startswith(tag):
            raw = raw[len(tag):]
    if raw.endswith("```"):
        raw = raw[:-3]

    try:
        return json.loads(raw.strip())
    except Exception as e:
        print(f"[JSON xatolik] {e}")
        return [{
            "title": "Xatolik yuz berdi",
            "content": "Javobni o'qib bo'lmadi",
            "bullets": [],
            "image_search_query": "error"
        }]

# ══════════════════════════════════════════════════
#  ASOSIY GENERATSIYA FUNKSIYASI
# ══════════════════════════════════════════════════

def generate_pptx(input_data: str, is_image: bool, output_name: str, slide_count: int) -> str:
    print(f"[PPTX] Ma'lumot tayyorlanmoqda ({slide_count} ta slayd)...")
    slides_data = _get_slides_from_gemini(input_data, is_image, slide_count)
    
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    black_clr = RGBColor(0, 0, 0)
    gray_clr  = RGBColor(80, 80, 80)

    # 1. Asosiy muqova slaydi
    cover_slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_white_bg(cover_slide)
    main_title = slides_data[0].get("title", "Fizika Taqdimoti") if slides_data else "Taqdimot"
    _text(cover_slide, main_title, 1.0, 2.8, 8.0, 2.0, size=44, bold=True, align=PP_ALIGN.CENTER)
    _text(cover_slide, "Toza, aniq format", 1.0, 4.2, 8.0, 0.8, size=18, color=gray_clr, align=PP_ALIGN.CENTER)

    # 2. Asosiy slaydlar
    for i, info in enumerate(slides_data, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_white_bg(slide)
        
        # Sarlavha
        _text(slide, info.get("title", f"Slayd {i}"), 0.5, 0.4, 9.0, 1.0, size=32, bold=True)
        
        # Rasm qo'yish (DuckDuckGo orqali)
        has_img = False
        keyword = info.get("image_keyword")
        if keyword:
            img_path = f"_ddgs_img_{i}.jpg"
            if _download_image(keyword, img_path):
                try:
                    # Rasmni o'ng tomonga joylashtiramiz
                    slide.shapes.add_picture(
                        img_path, Inches(5.8), Inches(1.8),
                        width=Inches(3.8) # balandlik rasmni o'ziga qarab moslashadi
                    )
                    has_img = True
                except Exception as e:
                    print(f"  [Rasm qoyish xatosi] {e}")
                finally:
                    if os.path.exists(img_path):
                        os.remove(img_path)

        text_width = 5.0 if has_img else 9.0

        # Asosiy matn
        content = info.get("content", "")
        if content:
            _text(slide, content, 0.5, 1.8, text_width, 2.0, size=16)

        # Bullets (Qo'shimcha nuqtalar)
        bullets = info.get("bullets", [])
        if bullets:
            bullet_text = "\n\n".join(f"- {b}" for b in bullets)
            _text(slide, bullet_text, 0.5, 3.5, text_width, 3.0, size=15, color=gray_clr)
            
        print(f"  -> Slayd {i} yasalmoqda: {info.get('title')}")

    filename = f"{output_name}.pptx"
    prs.save(filename)
    return filename
